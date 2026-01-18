import base64
import uuid
import requests
from bs4 import BeautifulSoup
from urllib.parse import urljoin, urlparse
from collections import deque
from typing import Dict, Any, List, Set, Tuple, Optional
import mimetypes
import io
import re
from app.config.settings import settings
from app.services.helper import _remove_fragment

# Optional imports for better extraction
HAVE_FITZ = False
try:
    import fitz  # PyMuPDF for robust PDF text + images
    HAVE_FITZ = True
except Exception:
    HAVE_FITZ = False

HAVE_DOCX = False
try:
    import docx  # python-docx
    HAVE_DOCX = True
except Exception:
    HAVE_DOCX = False

HAVE_PPTX = False
try:
    from pptx import Presentation  # python-pptx
    HAVE_PPTX = True
except Exception:
    HAVE_PPTX = False

HAVE_OPENPYXL = False
try:
    import openpyxl  # for .xlsx text extraction
    HAVE_OPENPYXL = True
except Exception:
    HAVE_OPENPYXL = False


class CrawlResult:
    def __init__(self):
        self.pages: Dict[str, Dict[str, Any]] = {}
        self.files: Dict[str, Dict[str, Any]] = {}
        self.images: Dict[str, Dict[str, Any]] = {}

def _registrable_domain(netloc: str) -> str:
    parts = netloc.split(".")
    return ".".join(parts[-2:]) if len(parts) >= 2 else netloc

def _same_domain_or_subdomain(target: str, base_reg_domain: str) -> bool:
    return target == base_reg_domain or target.endswith("." + base_reg_domain)

def _is_http_url(u: str) -> bool:
    p = urlparse(u)
    return p.scheme in ("http", "https")

def _normalize_url(base: str, href: str) -> str:
    return urljoin(base, href)

def _filename_from_url(url: str) -> Optional[str]:
    path = urlparse(url).path
    if not path:
        return None
    name = path.split("/")[-1]
    return name or None

def _split_large_text(text: str, max_chars: int = 15000) -> List[str]:
    """Split large text while preserving document structure"""
    if len(text) <= max_chars:
        return [text]
    
    chunks = []
    paragraphs = text.split('\n\n')
    current_chunk = []
    current_length = 0
    
    for para in paragraphs:
        para_length = len(para) + 2  # +2 for \n\n
        
        # If adding this paragraph would exceed chunk size
        if current_length + para_length > max_chars and current_chunk:
            # Finalize current chunk
            chunk_text = "\n\n".join(current_chunk).strip()
            if chunk_text:
                chunks.append(chunk_text)
            
            # Start new chunk with some overlap for context
            if len(current_chunk) > 2:
                # Keep last 2 paragraphs for context
                overlap = current_chunk[-2:]
                current_chunk = overlap + [para]
                current_length = sum(len(p) + 2 for p in overlap) + len(para)
            else:
                current_chunk = [para]
                current_length = para_length
        else:
            current_chunk.append(para)
            current_length += para_length
    
    # Handle final chunk
    if current_chunk:
        chunk_text = "\n\n".join(current_chunk).strip()
        if chunk_text:
            chunks.append(chunk_text)
    
    return chunks

def _fetch(url: str, headers: Dict[str, str]) -> Tuple[int, str, bytes, str]:
    r = requests.get(url, headers=headers, timeout=settings.REQUEST_TIMEOUT)
    ctype = (r.headers.get("Content-Type") or "").split(";")[0].strip().lower()
    return r.status_code, ctype, r.content, r.text

def _fetch_bytes(url: str, headers: Dict[str, str]) -> Tuple[bytes, str]:
    try:
        r = requests.get(url, headers=headers, timeout=settings.REQUEST_TIMEOUT)
        if r.status_code == 200:
            content_type = (r.headers.get("Content-Type") or "").split(";")[0].strip() \
                           or mimetypes.guess_type(url)[0] \
                           or "application/octet-stream"
            return r.content, content_type
    except Exception as e:
        print(f"Error fetching {url}: {e}")
    return b"", "application/octet-stream"

# ---------------- HTML Extraction ----------------

def _extract_visible_text_from_html(html: str) -> Tuple[str, Optional[str]]:
    soup = BeautifulSoup(html, "html.parser")

    # Remove script/style
    for bad in soup(["script", "style", "noscript"]):
        bad.extract()

    title = None
    if soup.title and soup.title.string:
        title = soup.title.string.strip()

    # Capture a wider set of elements that typically contain content
    text_parts = []
    for tag in soup.find_all(["h1","h2","h3","h4","h5","h6","p","li","blockquote","figcaption","pre","code","td","th"]):
        t = tag.get_text(" ", strip=True)
        if t:
            text_parts.append(t)

    # Fallback: if nothing captured, try body text
    if not text_parts and soup.body:
        all_text = soup.body.get_text(" ", strip=True)
        if all_text:
            text_parts.append(all_text)

    text = "\n".join(text_parts)
    # De-duplicate repeated spaces/newlines
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text, title

def _download_html_images(base_url: str, soup: BeautifulSoup, headers: Dict[str, str], downloaded: Set[str]) -> List[Dict[str, Any]]:
    images: List[Dict[str, Any]] = []
    inline_image_counter = 0
    
    for img in soup.find_all("img", src=True):
        src = _normalize_url(base_url, img["src"])

        if src.startswith("data:image"):
            # Handle base64 inline image
            header, b64data = src.split(",", 1)
            mime = header.split(":")[1].split(";")[0]
            img_bytes = base64.b64decode(b64data)
            
            # Create unique filename for each inline image
            inline_image_counter += 1
            filename = f"inline_image_{inline_image_counter}.{mime.split('/')[-1]}"

            images.append({
                "content": img_bytes,
                "mime": mime,
                "source": "inline",
                "filename": filename
            })
            continue
                
        if not _is_http_url(src):
            continue
        if src in downloaded:
            continue
        content, mime = _fetch_bytes(src, headers)
        if content:
            images.append({
                "content": content,
                "mime": mime,
                "source": src,
                "filename": _filename_from_url(src)
            })
            downloaded.add(src)
    return images

def _discover_and_download_linked_files(base_url: str, soup: BeautifulSoup, headers: Dict[str, str], downloaded: Set[str]) -> List[Dict[str, Any]]:
    files: List[Dict[str, Any]] = []
    exts = [".pdf", ".docx", ".pptx", ".xlsx", ".zip"]
    for a in soup.find_all("a", href=True):
        href = _normalize_url(base_url, a["href"])
        lower = href.lower()
        if any(lower.endswith(ext) for ext in exts):
            if not _is_http_url(href):
                continue
            if href in downloaded:
                continue
            content, mime = _fetch_bytes(href, headers)
            if content:
                files.append({
                    "content": content,
                    "mime": mime,
                    "source": href,
                    "filename": _filename_from_url(href)
                })
                downloaded.add(href)
    return files

# ---------------- Enhanced PDF Extraction with Layout Awareness ----------------

def _extract_pdf_layout_aware(raw: bytes, url: str = "") -> Tuple[str, List[Dict[str, Any]]]:
    """Extract PDF with layout awareness - Fixed version"""
    if not HAVE_FITZ:
        return "", []
        
    text_parts: List[str] = []
    images: List[Dict[str, Any]] = []

    try:
        with fitz.open(stream=raw, filetype="pdf") as doc:
            # Extract document name for unique image IDs
            doc_name = ""
            if url:
                from urllib.parse import urlparse
                parsed_url = urlparse(url)
                doc_name = parsed_url.path.split('/')[-1].replace('.pdf', '') if '.pdf' in parsed_url.path else "doc"
                doc_name = doc_name.replace('.', '_').replace('-', '_')
                if not doc_name:
                    doc_name = "document"
            
            for page_index in range(len(doc)):
                page = doc[page_index]
                
                # Get page content with layout information
                page_dict = page.get_text("dict")
                
                # Create layout grid
                layout_content = []
                image_counter = 0
                
                # Process all elements with their positions
                for block in page_dict.get("blocks", []):
                    bbox = block.get("bbox", [0,0,0,0])
                    x, y, x1, y1 = bbox
                    
                    if "lines" in block:  # Text
                        block_text = ""
                        for line in block["lines"]:
                            line_text = "".join(span["text"] for span in line.get("spans", []))
                            if line_text.strip():
                                block_text += line_text + " "
                        
                        if block_text.strip():
                            layout_content.append({
                                "type": "text",
                                "content": block_text.strip(),
                                "x": x,
                                "y": y,
                                "width": x1-x,
                                "height": y1-y,
                                "page": page_index + 1,
                                "bbox": bbox
                            })
                            
                    elif block.get("type", 0) == 1:  # Image
                        try:
                            # Try multiple methods to extract image
                            img_bytes = None
                            img_ext = "png"
                            mime = "image/png"
                            
                            # Method 1: Try xref-based extraction
                            xref = block.get("xref")
                            if xref is not None:
                                try:
                                    base = doc.extract_image(xref)
                                    img_bytes = base["image"]
                                    img_ext = base.get("ext", "png").lower()
                                    mime = f"image/{'jpeg' if img_ext in ('jpg','jpeg') else img_ext}"
                                    # Create globally unique image ID
                                    image_id = f"{doc_name}_page{page_index+1}_xref{xref}" if doc_name else f"page{page_index+1}_xref{xref}"
                                except Exception:
                                    pass
                            
                            # Method 2: If xref fails, try direct pixmap extraction
                            if img_bytes is None:
                                try:
                                    pix = page.get_pixmap(clip=bbox, matrix=fitz.Matrix(2, 2))
                                    img_bytes = pix.tobytes("png")
                                    img_ext = "png"
                                    mime = "image/png"
                                    # Create globally unique image ID
                                    image_id = f"{doc_name}_page{page_index+1}_area{int(x)}_{int(y)}" if doc_name else f"page{page_index+1}_area{int(x)}_{int(y)}"
                                except Exception:
                                    pass
                            
                            # Only add if we successfully extracted image data
                            if img_bytes is not None and len(img_bytes) > 0:
                                image_entry = {
                                    "content": img_bytes,
                                    "mime": mime,
                                    "x": x,
                                    "y": y,
                                    "width": x1-x,
                                    "height": y1-y,
                                    "page": page_index + 1,
                                    "source": image_id,
                                    "filename": f"{image_id}.{img_ext}",
                                    "bbox": bbox,
                                    "document_source": url  # Add document source for reference
                                }
                                
                                layout_content.append({
                                    "type": "image",
                                    "content": img_bytes,
                                    "mime": mime,
                                    "x": x,
                                    "y": y,
                                    "width": x1-x,
                                    "height": y1-y,
                                    "page": page_index + 1,
                                    "source": image_id,
                                    "filename": f"{image_id}.{img_ext}",
                                    "bbox": bbox
                                })
                                
                                images.append(image_entry)
                                image_counter += 1
                                print(f"Extracted image: {image_id} ({len(img_bytes)} bytes)")
                                
                        except Exception as e:
                            print(f"Image extraction skipped on page {page_index+1}: {e}")
                            continue
                
                # Sort by reading order
                layout_content.sort(key=lambda item: (item["page"], item["y"], item["x"]))
                
                # Reconstruct with proper flow
                page_text = f"[PAGE:{page_index+1}]\n"
                current_line_y = None
                line_items = []
                
                for item in layout_content:
                    if item["type"] == "text":
                        if current_line_y is None or abs(item["y"] - current_line_y) < 5:
                            line_items.append(item)
                            if current_line_y is None:
                                current_line_y = item["y"]
                        else:
                            line_items.sort(key=lambda x: x["x"])
                            line_text = " ".join(item["content"] for item in line_items)
                            page_text += line_text + "\n"
                            line_items = [item]
                            current_line_y = item["y"]
                    
                    elif item["type"] == "image":
                        # Use consistent format for image references - now with full context
                        page_text += f"[[IMAGE:{item['source']}]]\n"  # This will be like "fftester_page1_area12_18"
                
                # Process remaining line items
                if line_items:
                    line_items.sort(key=lambda x: x["x"])
                    line_text = " ".join(item["content"] for item in line_items)
                    page_text += line_text + "\n"
                
                page_text += f"[/PAGE:{page_index+1}]\n"
                text_parts.append(page_text)
        
        full_text = "\n".join(text_parts)
        
        # Add image summary at the end with document context
        if images:
            full_text += f"\n\n[DOCUMENT_IMAGES]\n"
            for i, img in enumerate(images, 1):
                full_text += f"Image {i}: [[IMAGE:{img['source']}]]\n"
        
        print(f"Successfully extracted {len(images)} images from document")
        return full_text, images
        
    except Exception as e:
        print(f"PDF extraction failed: {e}")
        import traceback
        traceback.print_exc()
        return _extract_pdf_basic(raw)

def _extract_pdf_basic(raw: bytes) -> Tuple[str, List[Dict[str, Any]]]:
    """Basic PDF extraction fallback with better error handling"""
    text_parts: List[str] = []
    images: List[Dict[str, Any]] = []
    
    if HAVE_FITZ:
        try:
            with fitz.open(stream=raw, filetype="pdf") as doc:
                for page_index in range(len(doc)):
                    page = doc[page_index]
                    text_parts.append(page.get_text("text"))
                    
                    # Extract images with comprehensive error handling
                    try:
                        page_images = page.get_images(full=True)
                        print(f"Found {len(page_images)} images on page {page_index+1}")
                        
                        for i, img in enumerate(page_images):
                            try:
                                xref = img[0]
                                print(f"Extracting image {i+1} with xref {xref}")
                                base = doc.extract_image(xref)
                                img_bytes = base["image"]
                                img_ext = base.get("ext", "png").lower()
                                mime = f"image/{'jpeg' if img_ext in ('jpg','jpeg') else img_ext}"
                                
                                if img_bytes and len(img_bytes) > 0:
                                    images.append({
                                        "content": img_bytes,
                                        "mime": mime,
                                        "source": f"page{page_index+1}_img{i+1}",
                                        "filename": f"page{page_index+1}_img{i+1}.{img_ext}"
                                    })
                                    print(f"Successfully extracted image {i+1}")
                                else:
                                    print(f"Image {i+1} has no data")
                                    
                            except Exception as img_error:
                                print(f"Failed to extract image {i+1}: {img_error}")
                                continue
                    except Exception as page_error:
                        print(f"Failed to get images from page {page_index+1}: {page_error}")
                        pass  # Continue with other pages
                        
        except Exception as e:
            print(f"Basic PDF extraction failed: {e}")
    
    text = "\n".join(p.strip() for p in text_parts if p and p.strip())
    print(f"Extracted {len(images)} images total")
    return text, images

# ---------------- DOCX / PPTX / XLSX (Optional) ----------------

def _extract_docx(raw: bytes) -> Tuple[str, List[Dict[str, Any]]]:
    if not HAVE_DOCX:
        return "", []
    text_parts: List[str] = []
    images: List[Dict[str, Any]] = []
    try:
        doc = docx.Document(io.BytesIO(raw))
        for p in doc.paragraphs:
            if p.text and p.text.strip():
                text_parts.append(p.text.strip())
        # Extract images via package parts
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                image_part = rel.target_part
                img_bytes = image_part.blob
                mime = getattr(image_part, "content_type", "image/png")
                images.append({
                    "content": img_bytes,
                    "mime": mime,
                    "source": "embedded:docx",
                    "filename": getattr(image_part, "partname", None)
                })
    except Exception as e:
        print(f"DOCX extraction failed: {e}")
    return "\n".join(text_parts), images

def _extract_pptx(raw: bytes) -> Tuple[str, List[Dict[str, Any]]]:
    if not HAVE_PPTX:
        return "", []
    text_parts: List[str] = []
    images: List[Dict[str, Any]] = []
    try:
        prs = Presentation(io.BytesIO(raw))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    text_parts.append(shape.text.strip())
                # Pictures
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        image = shape.image
                        img_bytes = image.blob
                        mime = image.content_type or "image/png"
                        images.append({
                            "content": img_bytes,
                            "mime": mime,
                            "source": "embedded:pptx",
                            "filename": getattr(image, "filename", None)
                        })
                    except Exception:
                        pass  # Some shapes might not be extractable
    except Exception as e:
        print(f"PPTX extraction failed: {e}")
    return "\n".join(text_parts), images

def _extract_xlsx(raw: bytes) -> Tuple[str, List[Dict[str, Any]]]:
    if not HAVE_OPENPYXL:
        return "", []
    text_parts: List[str] = []
    images: List[Dict[str, Any]] = []
    try:
        wb = openpyxl.load_workbook(io.BytesIO(raw), data_only=True)
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    text_parts.append(" | ".join(cells))
    except Exception as e:
        print(f"XLSX extraction failed: {e}")
    return "\n".join(text_parts), images

# ---------------- Router ----------------

def _extract_from_binary(raw: bytes, mime: str, url: str) -> Tuple[str, List[Dict[str, Any]]]:
    mime = (mime or "").lower()
    url_l = url.lower()

    if "pdf" in mime or url_l.endswith(".pdf"):
        return _extract_pdf_layout_aware(raw, url)  # Pass URL for better image naming
    elif "officedocument.wordprocessingml.document" in mime or url_l.endswith(".docx"):
        return _extract_docx(raw)
    elif "officedocument.presentationml.presentation" in mime or url_l.endswith(".pptx"):
        return _extract_pptx(raw)
    elif "officedocument.spreadsheetml.sheet" in mime or url_l.endswith(".xlsx"):
        return _extract_xlsx(raw)

    # Unsupported binary: no text/images extraction
    return "", []

# ---------------- Main crawl ----------------

async def crawl(root_url: str) -> CrawlResult:
    result = CrawlResult()
    headers = {"User-Agent": "Mozilla/5.0 (compatible; AI Assistant/1.0)"}

    root_parsed = urlparse(root_url)
    base_reg = _registrable_domain(root_parsed.netloc)

    queue = deque([(root_url, 0)])
    visited: Set[str] = set()
    downloaded_binary_sources: Set[str] = set()

    while queue and len(result.pages) < settings.MAX_PAGES:
        url, depth = queue.popleft()
        if url in visited or depth > settings.MAX_DEPTH:
            continue
        visited.add(url)

        try:
            status, content_type, raw, text_guess = _fetch(url, headers)
            if status != 200:
                print(f"Skip {url}: status {status}")
                continue

            doc_id = f"docs://{url}"

            # Non-HTML → treat as a document (e.g., PDF) and extract text + images
            if content_type and content_type != "text/html":
                try:
                    text, extracted_images = _extract_from_binary(raw, content_type, url)
                except Exception as e:
                    print(f"Binary extraction failed for {url}: {e}")
                    text, extracted_images = "", []
                
                # Store the complete document with all images first
                result.pages[doc_id] = {
                    "doc_id": doc_id,
                    "source_url": url,
                    "content_type": content_type,
                    "text": text,
                    "images": extracted_images,
                    "files": [{
                        "content": raw,
                        "mime": content_type or (mimetypes.guess_type(url)[0] or "application/octet-stream"),
                        "source": url,
                        "filename": _filename_from_url(url),
                    }],
                    "meta": {
                        "title": None, 
                        "status": status,
                        "word_count": len(text.split()),
                        "paragraph_count": len([p for p in text.split('\n\n') if p.strip()]),
                        "extraction_method": f"{content_type}_layout_aware",
                        "is_complete_document": True
                    }
                }
                
                print(f"Stored complete document {doc_id}: text_len={len(text)}, images={len(extracted_images)}")
                
                # Then create chunks if document is large
                max_doc_size = 20000
                if len(text) > max_doc_size:
                    text_chunks = _split_large_text(text, max_doc_size)
                    print(f"Splitting large document into {len(text_chunks)} chunks")
                    
                    # Create chunk references that point to the parent document
                    for i, chunk_text in enumerate(text_chunks):
                        chunk_doc_id = f"docs://{url}#chunk-{i+1}"
                        result.pages[chunk_doc_id] = {
                            "doc_id": chunk_doc_id,
                            "source_url": url,
                            "content_type": content_type,
                            "text": chunk_text,
                            "images": [],  # Empty - will look up from parent
                            "files": [],   # Empty - will look up from parent
                            "meta": {
                                "title": None,
                                "status": status,
                                "chunk_number": i+1,
                                "total_chunks": len(text_chunks),
                                "word_count": len(chunk_text.split()),
                                "extraction_method": f"{content_type}_layout_aware",
                                "parent_document": doc_id,
                                "is_chunk": True
                            }
                        }
                        print(f"Created chunk {i+1}/{len(text_chunks)} for {url}")
                
                # Continue to next URL since we've processed this binary document
                continue

            # HTML flow - Only process if we get here (content_type == "text/html" or content_type == None)
            soup = BeautifulSoup(text_guess, "html.parser")
            text, title = _extract_visible_text_from_html(text_guess)

            # Initialize images and files lists for HTML documents only
            images = []
            files = []
            
            try:
                # Download HTML images
                images = _download_html_images(url, soup, headers, downloaded_binary_sources)
                
                # Download linked documents
                files = _discover_and_download_linked_files(url, soup, headers, downloaded_binary_sources)

                # For any linked documents we just downloaded, attempt extraction too
                for f in files:
                    try:
                        f_text, f_imgs = _extract_from_binary(f["content"], f["mime"], f["source"])
                        # Append extracted elements to this page's content
                        if f_text:
                            text += ("\n\n[Content from linked document: " + (f["filename"] or f["source"]) + "]\n" + f_text)
                        # Add extracted images to the main images list
                        images.extend(f_imgs)
                    except Exception as e:
                        print(f"Extraction failed for linked file {f.get('source', 'unknown')}: {e}")
                        continue
                        
            except Exception as e:
                print(f"Error processing assets for {url}: {e}")

            # Store the complete HTML document with all images first
            result.pages[doc_id] = {
                "doc_id": doc_id,
                "source_url": url,
                "content_type": content_type,
                "text": text,
                "images": images,
                "files": files,
                "meta": {
                    "title": title,
                    "status": status,
                    "word_count": len(text.split()),
                    "paragraph_count": len([p for p in text.split('\n\n') if p.strip()]),
                    "extraction_method": "html_layout_aware",
                    "is_complete_document": True
                }
            }
            
            print(f"Stored complete HTML document {doc_id}: text_len={len(text)}, images={len(images)}, files={len(files)}")

            # Then create chunks if document is large
            max_doc_size = 20000
            if len(text) > max_doc_size:
                text_chunks = _split_large_text(text, max_doc_size)
                print(f"Splitting large HTML document into {len(text_chunks)} chunks")
                
                # Create chunk references that point to the parent document
                for i, chunk_text in enumerate(text_chunks):
                    chunk_doc_id = f"{doc_id}#chunk-{i+1}"
                    result.pages[chunk_doc_id] = {
                        "doc_id": chunk_doc_id,
                        "source_url": url,
                        "content_type": content_type,
                        "text": chunk_text,
                        "images": [],  # Empty - will look up from parent
                        "files": [],   # Empty - will look up from parent
                        "meta": {
                            "title": title,
                            "status": status,
                            "chunk_number": i+1,
                            "total_chunks": len(text_chunks),
                            "word_count": len(chunk_text.split()),
                            "paragraph_count": len([p for p in chunk_text.split('\n\n') if p.strip()]),
                            "extraction_method": "html_layout_aware",
                            "parent_document": doc_id,
                            "is_chunk": True
                        }
                    }
                    print(f"Created HTML chunk {i+1}/{len(text_chunks)} for {url}")

            # Enqueue internal links
            for a in soup.find_all("a", href=True):
                link = _normalize_url(url, a["href"])
                link = _remove_fragment(link)
                if not _is_http_url(link):
                    continue
                parsed = urlparse(link)
                if _same_domain_or_subdomain(parsed.netloc, base_reg):
                    if link not in visited and len(result.pages) + len(queue) < settings.MAX_PAGES:
                        queue.append((link, depth + 1))

        except Exception as e:
            print(f"Error crawling {url}: {e}")
            import traceback
            traceback.print_exc()
    
    return result